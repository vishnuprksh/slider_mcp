"""PPTX Rendering Engine.

Converts a DeckSpec into a .pptx file using python-pptx.

Slide dimensions, color tokens, and content blocks are translated
faithfully from the DeckSpec schema. All rendering is synchronous.

python-pptx coordinate system: EMU (English Metric Units)
  1 inch = 914400 EMU
  Standard 16:9 slide = 12192000 x 6858000 EMU
"""
from __future__ import annotations

import io
import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

from app.models.deck import (
    AspectRatio,
    BulletListBlock,
    CodeBlock,
    ContentBlockType,
    DeckSpec,
    IconBlock,
    ImageBlock,
    QuoteBlock,
    SlideLayout,
    SlideSpec,
    StatBlock,
    SVGBlock,
    TextBlock,
    ThemeSpec,
)


# ─────────────────────────────────────────────────────────────────────────────
# Slide dimensions (EMU)
# ─────────────────────────────────────────────────────────────────────────────

_SLIDE_DIMS: dict[AspectRatio, tuple[int, int]] = {
    AspectRatio.WIDESCREEN: (12192000, 6858000),   # 16:9
    AspectRatio.STANDARD: (9144000, 6858000),       # 4:3
    AspectRatio.SQUARE: (6858000, 6858000),          # 1:1
    AspectRatio.PORTRAIT: (6858000, 12192000),       # 9:16
}


def _slide_dims(ratio: AspectRatio) -> tuple[int, int]:
    return _SLIDE_DIMS.get(ratio, _SLIDE_DIMS[AspectRatio.WIDESCREEN])


# ─────────────────────────────────────────────────────────────────────────────
# Color helpers
# ─────────────────────────────────────────────────────────────────────────────


def _parse_hex(color: str) -> RGBColor:
    """Parse a #RRGGBB or #RGB hex string to RGBColor.
    Falls back to black on invalid input."""
    color = color.strip().lstrip("#")
    if len(color) == 3:
        color = "".join(c * 2 for c in color)
    try:
        r, g, b = int(color[0:2], 16), int(color[2:4], 16), int(color[4:6], 16)
        return RGBColor(r, g, b)
    except (ValueError, IndexError):
        return RGBColor(0, 0, 0)


def _is_hex_color(color: str) -> bool:
    return bool(re.match(r"^#?[0-9a-fA-F]{3,8}$", color.strip()))


# ─────────────────────────────────────────────────────────────────────────────
# Text style mapping
# ─────────────────────────────────────────────────────────────────────────────

_TEXT_STYLE_PT: dict[str, int] = {
    "h1": 40, "h2": 32, "h3": 24, "h4": 20,
    "body": 16, "caption": 12, "label": 10,
}

_ALIGN_MAP: dict[str, PP_ALIGN] = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
}


# ─────────────────────────────────────────────────────────────────────────────
# Slide-level helpers
# ─────────────────────────────────────────────────────────────────────────────


def _blank_layout(prs: Presentation):
    """Return the blank slide layout from the presentation."""
    # Layout index 6 is 'Blank' in the default template; fall back gracefully.
    try:
        return prs.slide_layouts[6]
    except IndexError:
        return prs.slide_layouts[-1]


def _fill_slide_bg(slide, color: RGBColor) -> None:
    """Set the slide background to a solid color."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_textbox(slide, left: int, top: int, width: int, height: int,
                 text: str, font_size: int = 16, bold: bool = False,
                 color: RGBColor | None = None,
                 align: PP_ALIGN = PP_ALIGN.LEFT,
                 wrap: bool = True) -> None:
    """Add a text box to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color


# ─────────────────────────────────────────────────────────────────────────────
# Content block renderers
# ─────────────────────────────────────────────────────────────────────────────

# Layout constants (EMU)
_MARGIN = Inches(0.6)


def _render_block_to_slide(slide, block, theme: ThemeSpec,
                            left: int, top: int, width: int, height: int,
                            y_cursor: list[int]) -> int:
    """Render a single content block onto a slide. Returns new y_cursor."""
    bt = block.block_type
    text_color = _parse_hex(theme.colors.text_primary) if _is_hex_color(theme.colors.text_primary) else RGBColor(0x0F, 0x17, 0x2A)
    muted_color = _parse_hex(theme.colors.text_secondary) if _is_hex_color(theme.colors.text_secondary) else RGBColor(0x47, 0x55, 0x69)
    primary_color = _parse_hex(theme.colors.primary) if _is_hex_color(theme.colors.primary) else RGBColor(0x25, 0x63, 0xEB)
    accent_color = _parse_hex(theme.colors.accent) if _is_hex_color(theme.colors.accent) else RGBColor(0xF5, 0x9E, 0x0B)

    y = y_cursor[0]
    block_height = Inches(0.4)  # default

    if bt == ContentBlockType.TEXT:
        fs = _TEXT_STYLE_PT.get(block.style, 16)
        bold = block.style in {"h1", "h2", "h3", "h4"}
        align = _ALIGN_MAP.get(block.align, PP_ALIGN.LEFT)
        block_height = Pt(fs) * 2
        color = text_color
        if block.color and _is_hex_color(block.color):
            color = _parse_hex(block.color)
        _add_textbox(slide, left, y, width, int(block_height),
                     block.text, fs, bold, color, align)
        y += int(block_height) + Inches(0.1)

    elif bt == ContentBlockType.BULLET_LIST:
        txBox = slide.shapes.add_textbox(left, y, width, int(Inches(0.35) * len(block.items)))
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, item in enumerate(block.items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            prefix = "→ " if not block.ordered else f"{i+1}. "
            run = p.add_run()
            run.text = f"{prefix}{item}"
            run.font.size = Pt(16)
            if i == 0 and block.highlight_first:
                run.font.bold = True
                run.font.color.rgb = primary_color
            else:
                run.font.color.rgb = text_color
        item_height = int(Inches(0.35) * len(block.items))
        y += item_height + Inches(0.15)

    elif bt == ContentBlockType.STAT:
        # Value (large)
        _add_textbox(slide, left, y, width, int(Inches(0.8)),
                     block.value, 40, True, primary_color, PP_ALIGN.CENTER)
        # Label (small)
        _add_textbox(slide, left, y + int(Inches(0.8)), width, int(Inches(0.35)),
                     block.label, 12, False, muted_color, PP_ALIGN.CENTER)
        y += int(Inches(1.25))

    elif bt == ContentBlockType.QUOTE:
        quote_text = f'"{block.text}"'
        _add_textbox(slide, left, y, width, int(Inches(0.8)),
                     quote_text, 18, False, text_color, PP_ALIGN.CENTER)
        if block.attribution:
            attr = f"— {block.attribution}"
            if block.role:
                attr += f", {block.role}"
            _add_textbox(slide, left, y + int(Inches(0.85)), width, int(Inches(0.35)),
                         attr, 12, False, muted_color, PP_ALIGN.CENTER)
        y += int(Inches(1.3))

    elif bt == ContentBlockType.CODE:
        _add_textbox(slide, left, y, width, int(Inches(1.2)),
                     block.code, 11, False, RGBColor(0xCD, 0xD6, 0xF4))
        y += int(Inches(1.3))

    elif bt == ContentBlockType.IMAGE:
        # Placeholder text (embedding arbitrary URLs in pptx is unsafe)
        try:
            ph = slide.shapes.add_textbox(left, y, width, int(Inches(1.5)))
            tf = ph.text_frame
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = f"[Image: {block.asset.alt_text or block.asset.source[:40]}]"
            run.font.size = Pt(12)
            run.font.color.rgb = muted_color
        except Exception:
            pass
        y += int(Inches(1.6))

    elif bt == ContentBlockType.ICON:
        _add_textbox(slide, left, y, width, int(Inches(0.4)),
                     f"[Icon: {block.icon.name}]", 12, False, muted_color, PP_ALIGN.CENTER)
        if block.label:
            _add_textbox(slide, left, y + int(Inches(0.4)), width, int(Inches(0.3)),
                         block.label, 10, False, muted_color, PP_ALIGN.CENTER)
        y += int(Inches(0.8))

    elif bt == ContentBlockType.SVG:
        _add_textbox(slide, left, y, width, int(Inches(0.4)),
                     "[SVG graphic]", 10, False, muted_color, PP_ALIGN.CENTER)
        y += int(Inches(0.5))

    y_cursor[0] = y
    return y


# ─────────────────────────────────────────────────────────────────────────────
# Slide builder
# ─────────────────────────────────────────────────────────────────────────────


def _render_slide(prs: Presentation, slide_spec: SlideSpec, theme: ThemeSpec,
                  slide_w: int, slide_h: int) -> None:
    """Add one slide to the presentation."""
    layout = _blank_layout(prs)
    slide = prs.slides.add_slide(layout)

    primary_rgb = _parse_hex(theme.colors.primary) if _is_hex_color(theme.colors.primary) else RGBColor(0x25, 0x63, 0xEB)
    bg_rgb = _parse_hex(theme.colors.background) if _is_hex_color(theme.colors.background) else RGBColor(0xFF, 0xFF, 0xFF)
    on_primary_rgb = _parse_hex(theme.colors.text_on_primary) if _is_hex_color(theme.colors.text_on_primary) else RGBColor(0xFF, 0xFF, 0xFF)
    text_rgb = _parse_hex(theme.colors.text_primary) if _is_hex_color(theme.colors.text_primary) else RGBColor(0x0F, 0x17, 0x2A)

    ltype = slide_spec.layout.layout

    # ── Full-bleed layouts (title, section_break, closing) ──────────────────
    if ltype in {SlideLayout.TITLE, SlideLayout.SECTION_BREAK, SlideLayout.CLOSING}:
        _fill_slide_bg(slide, primary_rgb)
        title_top = int(slide_h * 0.3)
        _add_textbox(slide, _MARGIN, title_top, slide_w - 2 * _MARGIN, int(Inches(1.2)),
                     slide_spec.title, 40, True, on_primary_rgb, PP_ALIGN.CENTER)
        if slide_spec.subtitle:
            _add_textbox(slide, _MARGIN, title_top + int(Inches(1.3)),
                         slide_w - 2 * _MARGIN, int(Inches(0.6)),
                         slide_spec.subtitle, 20, False, on_primary_rgb, PP_ALIGN.CENTER)
        # Content blocks below subtitle
        y_cursor = [title_top + int(Inches(2.2))]
        for block in slide_spec.content_blocks:
            _render_block_to_slide(slide, block, theme, _MARGIN, 0,
                                   slide_w - 2 * _MARGIN, slide_h, y_cursor)

    # ── Multi-column layouts ─────────────────────────────────────────────────
    elif ltype in {SlideLayout.TWO_COLUMN, SlideLayout.THREE_COLUMN, SlideLayout.COMPARISON}:
        _fill_slide_bg(slide, bg_rgb)
        header_h = int(Inches(1.0))
        _add_textbox(slide, _MARGIN, _MARGIN, slide_w - 2 * _MARGIN, header_h,
                     slide_spec.title, 28, True, text_rgb)
        # Distribute blocks across columns
        ratios = slide_spec.layout.column_ratios
        n_cols = len(ratios)
        total_ratio = sum(ratios)
        content_blocks = slide_spec.content_blocks
        per_col = max(1, (len(content_blocks) + n_cols - 1) // n_cols)
        col_x = _MARGIN
        content_top = _MARGIN + header_h + int(Inches(0.2))
        content_h = slide_h - content_top - _MARGIN
        available_w = slide_w - 2 * _MARGIN
        for col_idx in range(n_cols):
            col_ratio = ratios[col_idx] if col_idx < len(ratios) else 1
            col_w = int(available_w * col_ratio / total_ratio)
            col_blocks = content_blocks[col_idx * per_col:(col_idx + 1) * per_col]
            y_cursor = [content_top]
            for block in col_blocks:
                _render_block_to_slide(slide, block, theme, col_x, content_top,
                                       col_w, content_h, y_cursor)
            col_x += col_w + int(Inches(0.1))

    # ── Standard single-column ───────────────────────────────────────────────
    else:
        _fill_slide_bg(slide, bg_rgb)
        header_h = int(Inches(1.0))
        _add_textbox(slide, _MARGIN, _MARGIN, slide_w - 2 * _MARGIN, header_h,
                     slide_spec.title, 28, True, text_rgb)
        content_top = _MARGIN + header_h + int(Inches(0.2))
        y_cursor = [content_top]
        for block in slide_spec.content_blocks:
            _render_block_to_slide(slide, block, theme, _MARGIN, content_top,
                                   slide_w - 2 * _MARGIN, slide_h - content_top - _MARGIN,
                                   y_cursor)


# ─────────────────────────────────────────────────────────────────────────────
# Public API
# ─────────────────────────────────────────────────────────────────────────────


def render_pptx(deck: DeckSpec) -> bytes:
    """Render a DeckSpec to PPTX bytes.

    Args:
        deck: A schema-valid DeckSpec.

    Returns:
        Raw .pptx file content as bytes.
    """
    slide_w, slide_h = _slide_dims(deck.aspect_ratio)

    prs = Presentation()
    prs.slide_width = Emu(slide_w)
    prs.slide_height = Emu(slide_h)

    for slide_spec in deck.slides:
        _render_slide(prs, slide_spec, deck.theme, slide_w, slide_h)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def render_pptx_to_file(deck: DeckSpec, output_path: Path) -> Path:
    """Render deck to a .pptx file. Creates parent directories if needed.

    Args:
        deck: A schema-valid DeckSpec.
        output_path: Destination .pptx file path.

    Returns:
        Resolved output path.
    """
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    data = render_pptx(deck)
    output_path.write_bytes(data)
    return output_path.resolve()
