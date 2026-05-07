"""Phase 6 — PPTX renderer tests."""
from __future__ import annotations

import io
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Emu

from app.models.deck import (
    AspectRatio,
    BulletListBlock,
    CodeBlock,
    DeckSpec,
    LayoutSpec,
    QuoteBlock,
    SlideLayout,
    SlideRole,
    SlideSpec,
    StatBlock,
    TextBlock,
    ThemeSpec,
)
from app.renderers.pptx_renderer import _parse_hex, render_pptx, render_pptx_to_file
from app.services.planning import PlanningRequest, plan_deck


# ─────────────────────────────────────────────────────────────────────────────
# Helpers
# ─────────────────────────────────────────────────────────────────────────────


def load_prs(data: bytes) -> Presentation:
    return Presentation(io.BytesIO(data))


def minimal_deck(n: int = 3) -> DeckSpec:
    return plan_deck(PlanningRequest(topic="PPTX Test Deck", slide_count=n))


def slide_with(title: str, layout: SlideLayout, *blocks) -> SlideSpec:
    return SlideSpec(
        title=title,
        role=SlideRole.CONTENT,
        layout=LayoutSpec(layout=layout),
        content_blocks=list(blocks),
    )


# ─────────────────────────────────────────────────────────────────────────────
# _parse_hex
# ─────────────────────────────────────────────────────────────────────────────


class TestParseHex:
    def test_6char_hex(self):
        color = _parse_hex("#2563EB")
        assert color[0] == 0x25
        assert color[1] == 0x63
        assert color[2] == 0xEB

    def test_3char_hex_expanded(self):
        color = _parse_hex("#FFF")
        assert color == (0xFF, 0xFF, 0xFF)

    def test_no_hash_prefix(self):
        color = _parse_hex("FF0000")
        assert color[0] == 0xFF

    def test_invalid_falls_back_to_black(self):
        color = _parse_hex("not-a-color")
        assert color == (0, 0, 0)


# ─────────────────────────────────────────────────────────────────────────────
# render_pptx — format correctness
# ─────────────────────────────────────────────────────────────────────────────


class TestRenderPptx:
    def test_returns_bytes(self):
        deck = minimal_deck()
        result = render_pptx(deck)
        assert isinstance(result, bytes)

    def test_valid_pptx_magic_bytes(self):
        # PPTX is a ZIP; starts with PK magic bytes
        data = render_pptx(minimal_deck())
        assert data[:2] == b"PK"

    def test_slide_count_matches(self):
        deck = minimal_deck(5)
        prs = load_prs(render_pptx(deck))
        assert len(prs.slides) == 5

    def test_widescreen_dimensions(self):
        deck = minimal_deck()
        deck = deck.model_copy(update={"aspect_ratio": AspectRatio.WIDESCREEN})
        prs = load_prs(render_pptx(deck))
        assert prs.slide_width == Emu(12192000)
        assert prs.slide_height == Emu(6858000)

    def test_standard_4_3_dimensions(self):
        deck = minimal_deck()
        deck = deck.model_copy(update={"aspect_ratio": AspectRatio.STANDARD})
        prs = load_prs(render_pptx(deck))
        assert prs.slide_width == Emu(9144000)

    def test_portrait_dimensions(self):
        deck = minimal_deck()
        deck = deck.model_copy(update={"aspect_ratio": AspectRatio.PORTRAIT})
        prs = load_prs(render_pptx(deck))
        assert prs.slide_width == Emu(6858000)
        assert prs.slide_height == Emu(12192000)

    def test_title_appears_in_slide_shapes(self):
        deck = minimal_deck(3)
        prs = load_prs(render_pptx(deck))
        first_slide = prs.slides[0]
        texts = [
            run.text
            for shape in first_slide.shapes
            if shape.has_text_frame
            for para in shape.text_frame.paragraphs
            for run in para.runs
        ]
        assert deck.slides[0].title in " ".join(texts)

    def test_non_empty_file(self):
        data = render_pptx(minimal_deck(3))
        assert len(data) > 5000  # non-trivial PPTX


# ─────────────────────────────────────────────────────────────────────────────
# Content block rendering
# ─────────────────────────────────────────────────────────────────────────────


class TestPptxContentBlocks:
    def _render_deck(self, *slides: SlideSpec) -> Presentation:
        deck = DeckSpec(title="Block Test", slides=list(slides))
        return load_prs(render_pptx(deck))

    def test_text_block_text_present(self):
        slide = slide_with("T", SlideLayout.BULLETS, TextBlock(text="Unique Text Content XYZ"))
        prs = self._render_deck(slide)
        all_text = " ".join(
            run.text
            for shape in prs.slides[0].shapes
            if shape.has_text_frame
            for para in shape.text_frame.paragraphs
            for run in para.runs
        )
        assert "Unique Text Content XYZ" in all_text

    def test_bullet_list_items_present(self):
        slide = slide_with("T", SlideLayout.BULLETS,
                           BulletListBlock(items=["Alpha Item", "Beta Item"]))
        prs = self._render_deck(slide)
        all_text = " ".join(
            run.text
            for shape in prs.slides[0].shapes
            if shape.has_text_frame
            for para in shape.text_frame.paragraphs
            for run in para.runs
        )
        assert "Alpha Item" in all_text

    def test_stat_block_value_present(self):
        slide = slide_with("Stats", SlideLayout.STATS, StatBlock(value="99%", label="Uptime"))
        prs = self._render_deck(slide)
        all_text = " ".join(
            run.text
            for shape in prs.slides[0].shapes
            if shape.has_text_frame
            for para in shape.text_frame.paragraphs
            for run in para.runs
        )
        assert "99%" in all_text

    def test_quote_block_text_present(self):
        slide = slide_with("Q", SlideLayout.QUOTE,
                           QuoteBlock(text="The future is now", attribution="Futurist"))
        prs = self._render_deck(slide)
        all_text = " ".join(
            run.text
            for shape in prs.slides[0].shapes
            if shape.has_text_frame
            for para in shape.text_frame.paragraphs
            for run in para.runs
        )
        assert "The future is now" in all_text


# ─────────────────────────────────────────────────────────────────────────────
# Layout-specific rendering
# ─────────────────────────────────────────────────────────────────────────────


class TestPptxLayouts:
    def test_title_layout_renders(self):
        slide = SlideSpec(
            title="Big Title",
            role=SlideRole.COVER,
            layout=LayoutSpec(layout=SlideLayout.TITLE),
            content_blocks=[],
        )
        deck = DeckSpec(title="T", slides=[slide])
        data = render_pptx(deck)
        prs = load_prs(data)
        assert len(prs.slides) == 1

    def test_two_column_layout_renders(self):
        slide = SlideSpec(
            title="Two Col",
            role=SlideRole.CONTENT,
            layout=LayoutSpec(layout=SlideLayout.TWO_COLUMN, column_ratios=[50, 50]),
            content_blocks=[TextBlock(text="Left side"), TextBlock(text="Right side")],
        )
        deck = DeckSpec(title="T", slides=[slide])
        data = render_pptx(deck)
        assert load_prs(data) is not None

    def test_closing_layout_renders(self):
        slide = SlideSpec(
            title="Thank You",
            role=SlideRole.CLOSING,
            layout=LayoutSpec(layout=SlideLayout.CLOSING),
            content_blocks=[],
        )
        deck = DeckSpec(title="T", slides=[slide])
        data = render_pptx(deck)
        assert load_prs(data) is not None


# ─────────────────────────────────────────────────────────────────────────────
# render_pptx_to_file
# ─────────────────────────────────────────────────────────────────────────────


class TestRenderPptxToFile:
    def test_creates_file(self, tmp_path):
        deck = minimal_deck(3)
        out = render_pptx_to_file(deck, tmp_path / "output.pptx")
        assert out.exists()

    def test_file_is_valid_pptx(self, tmp_path):
        deck = minimal_deck(3)
        out = render_pptx_to_file(deck, tmp_path / "deck.pptx")
        data = out.read_bytes()
        assert data[:2] == b"PK"

    def test_creates_parent_dirs(self, tmp_path):
        deck = minimal_deck(3)
        nested = tmp_path / "a" / "b" / "deck.pptx"
        out = render_pptx_to_file(deck, nested)
        assert out.exists()

    def test_end_to_end_plan_to_pptx(self, tmp_path):
        req = PlanningRequest(topic="Data Science Overview", audience="engineering", slide_count=8)
        deck = plan_deck(req)
        out = render_pptx_to_file(deck, tmp_path / "ds.pptx")
        prs = load_prs(out.read_bytes())
        assert len(prs.slides) == 8
